import os
import io
import docx2txt
import PyPDF2
from docx import Document
from pptx import Presentation
import zipfile
import tempfile
import mimetypes
import magic  # For better file type detection

def extract_text_from_docx(file_content):
    """
    Extract text from a .docx file
    """
    try:
        # First try with docx2txt which preserves some formatting
        text = docx2txt.process(io.BytesIO(file_content))
        return text
    except Exception as e:
        try:
            # Fallback to python-docx
            doc = Document(io.BytesIO(file_content))
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return '\n'.join(full_text)
        except Exception as e:
            return f"Error extracting text from Word document: {str(e)}"

def extract_text_from_pdf(file_content):
    """
    Extract text from a PDF file
    """
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            text += pdf_reader.pages[page_num].extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error extracting text from PDF: {str(e)}"

def extract_text_from_pptx(file_content):
    """
    Extract text from a PowerPoint file
    """
    try:
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting text from PowerPoint: {str(e)}"

def is_binary_file(file_content, filename):
    """
    Check if a file is binary (non-text) based on content and extension
    """
    # First check by extension
    text_extensions = [
        '.txt', '.py', '.js', '.html', '.css', '.md', '.json', '.xml', 
        '.csv', '.c', '.cpp', '.h', '.java', '.php', '.rb', '.go', 
        '.rs', '.ts', '.sh', '.bat', '.ps1', '.yaml', '.yml', 
        '.toml', '.ini', '.cfg', '.docx', '.doc', '.pdf', '.pptx', '.ppt'
    ]
    
    ext = os.path.splitext(filename)[1].lower()
    if ext not in text_extensions:
        # Try to detect file type using python-magic
        try:
            mime = magic.Magic(mime=True)
            file_type = mime.from_buffer(file_content)
            
            # Check if it's a text type
            if not file_type.startswith('text/') and not any(x in file_type for x in ['word', 'pdf', 'powerpoint', 'excel', 'opendocument']):
                return True  # It's binary
        except:
            # If magic fails, use a simple heuristic
            try:
                content_sample = file_content[:1024]  # Check first 1KB
                return b'\0' in content_sample  # If it contains null bytes, it's likely binary
            except:
                return True  # Assume binary if we can't determine
    
    return False  # Assume it's a text file if we get here

def extract_text_from_file(file_content, filename):
    """
    Extract text from various file types
    """
    ext = os.path.splitext(filename)[1].lower()
    
    # Check if it's a binary file we don't support
    if is_binary_file(file_content, filename) and ext not in ['.docx', '.pdf', '.pptx', '.doc', '.ppt']:
        return None  # Skip binary files we don't support
    
    # Handle different file types
    if ext == '.docx' or ext == '.doc':
        return extract_text_from_docx(file_content)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_content)
    elif ext == '.pptx' or ext == '.ppt':
        return extract_text_from_pptx(file_content)
    elif ext == '.zip':
        # For zip files, we'll extract and process them separately
        return None
    else:
        # For text files, try to decode
        try:
            return file_content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                return file_content.decode('latin-1')  # Try another encoding
            except:
                return None  # Skip if we can't decode

def process_zip_file(zip_content, target_dict, max_depth=3, current_depth=0):
    """
    Process a zip file and extract text from supported files
    Handles nested zip files up to max_depth
    """
    if current_depth >= max_depth:
        return target_dict  # Prevent too deep recursion
    
    with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as temp_zip:
        temp_zip.write(zip_content)
        temp_zip_path = temp_zip.name
    
    try:
        with zipfile.ZipFile(temp_zip_path, 'r') as zip_ref:
            for file_info in zip_ref.infolist():
                if file_info.is_dir():
                    continue
                
                with zip_ref.open(file_info) as file:
                    content = file.read()
                    
                    # Handle nested zip files
                    if file_info.filename.lower().endswith('.zip'):
                        if current_depth < max_depth - 1:  # Process nested zip if not too deep
                            nested_prefix = f"[{current_depth+1}]_"
                            process_zip_file(content, target_dict, max_depth, current_depth + 1)
                        continue
                    
                    # Extract text from the file
                    extracted_text = extract_text_from_file(content, file_info.filename)
                    if extracted_text:
                        target_dict[file_info.filename] = extracted_text
    except Exception as e:
        print(f"Error processing zip file: {str(e)}")
    
    os.unlink(temp_zip_path)
    return target_dict

def is_resume(filename, content):
    """
    Attempt to detect if a file is likely a resume
    """
    # Check filename for resume-related terms
    resume_terms = ['resume', 'cv', 'curriculum', 'vitae']
    if any(term in filename.lower() for term in resume_terms):
        return True
    
    # Check content for resume-related sections
    resume_sections = [
        'experience', 'education', 'skills', 'objective', 
        'summary', 'employment', 'work history', 'qualification',
        'certification', 'reference', 'project'
    ]
    
    # For text content
    if isinstance(content, str):
        if any(section.lower() in content.lower() for section in resume_sections):
            # Count how many resume sections are present
            section_count = sum(1 for section in resume_sections if section.lower() in content.lower())
            # If multiple resume sections are present, it's likely a resume
            return section_count >= 3
    
    return False
